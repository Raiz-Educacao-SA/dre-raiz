from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BLUE  = RGBColor(0x15, 0x2e, 0x55)
MID_BLUE   = RGBColor(0x1B, 0x75, 0xBB)
ORANGE     = RGBColor(0xF4, 0x4C, 0x00)
CYAN       = RGBColor(0x4A, 0xC8, 0xF4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
GREEN      = RGBColor(0x05, 0x7A, 0x55)
PURPLE     = RGBColor(0x5B, 0x21, 0xB6)
HEADER_BLUE = RGBColor(0x1B, 0x3A, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def txtml(slide, lines, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def header(slide, title):
    rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
    rect(slide, 0, 0, 13.33, 1.1, DARK_BLUE)
    rect(slide, 0, 1.08, 13.33, 0.04, CYAN)
    rect(slide, 0, 1.1, 0.08, 6.4, ORANGE)
    txt(slide, title, 0.3, 0.18, 12, 0.75, 24, bold=True, color=WHITE)

def dark_header(slide, title):
    rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
    rect(slide, 0, 0, 13.33, 1.1, MID_BLUE)
    rect(slide, 0, 1.08, 13.33, 0.04, CYAN)
    rect(slide, 0, 1.1, 0.08, 6.4, ORANGE)
    txt(slide, title, 0.3, 0.18, 12, 0.75, 24, bold=True, color=WHITE)

# ── SLIDE 1: Capa ────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, 13.33, 7.5, DARK_BLUE)
rect(s1, 0, 5.8, 13.33, 1.7, MID_BLUE)
rect(s1, 0, 5.75, 13.33, 0.08, CYAN)
rect(s1, 0.5, 1.2, 0.08, 4.2, ORANGE)
txt(s1, "RAIZ EDUCACAO", 0.75, 1.1, 10, 0.7, 13, bold=True, color=CYAN)
txt(s1, "DRE Gerencial", 0.75, 1.8, 11, 1.5, 52, bold=True, color=WHITE)
txt(s1, "Plataforma de Gestao Financeira", 0.75, 3.25, 10, 0.65, 22, color=CYAN)
txt(s1, "Analise  .  Controle  .  Decisao", 0.75, 3.9, 8, 0.5, 14, color=RGBColor(0xB0, 0xC8, 0xE0))
txt(s1, "Confidencial  |  Raiz Educacao S.A.  |  2026", 0.75, 6.1, 12, 0.4, 10, color=RGBColor(0xCC, 0xDD, 0xEE))

# ── SLIDE 2: O que e a plataforma ────────────────────────────────
s2 = prs.slides.add_slide(blank)
header(s2, "O que e a Plataforma?")

rect(s2, 0.3, 1.35, 5.9, 5.55, WHITE)
rect(s2, 0.3, 1.35, 5.9, 0.08, MID_BLUE)
txt(s2, "VISAO GERAL", 0.55, 1.5, 5.4, 0.4, 11, bold=True, color=MID_BLUE)
txtml(s2, [
    "A plataforma DRE Gerencial e um sistema web desenvolvido",
    "exclusivamente para a Raiz Educacao, com o objetivo de",
    "centralizar e simplificar a gestao financeira da empresa.",
    "",
    "Integrada ao banco de dados Supabase e protegida por",
    "autenticacao Firebase, a ferramenta consolida dados de",
    "receitas, custos e despesas em tempo real, permitindo",
    "analises gerenciais por cenario, marca, filial e categoria.",
    "",
    "Acessivel por qualquer dispositivo via navegador, sem",
    "necessidade de instalacao ou configuracao local.",
], 0.55, 2.0, 5.5, 4.7, 11, color=DARK_TEXT)

rect(s2, 6.9, 1.35, 6.1, 5.55, WHITE)
rect(s2, 6.9, 1.35, 6.1, 0.08, ORANGE)
txt(s2, "PRINCIPAIS CARACTERISTICAS", 7.15, 1.5, 5.7, 0.4, 11, bold=True, color=ORANGE)
txtml(s2, [
    "  100% web  -  acesso via dre-raiz.vercel.app",
    "  Dados em tempo real direto do Supabase",
    "  Controle de acesso por perfil de usuario",
    "  Admin  /  Aprovador  /  Usuario padrao",
    "  Visao consolidada ou por cenario",
    "  Real, Orcado e Ano Anterior (A-1)",
    "  Drill-down por marca, filial, fornecedor e tag",
    "  Exportacao para Excel com formatacao profissional",
    "  Fluxo de aprovacao de ajustes com rastreabilidade",
    "  Rateio automatico de custos administrativos",
    "  Dashboard com KPIs executivos",
], 7.15, 2.0, 5.75, 4.7, 11, color=DARK_TEXT)

# ── SLIDE 3: Modulos ─────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
header(s3, "Modulos da Plataforma")

modules = [
    ("DRE Gerencial",  MID_BLUE,
     ["Demonstrativo de resultado com 3 modos:", "Consolidado, por Cenario e por Mes.",
      "Drill-down por marca, filial, fornecedor", "e tags. Export Excel fiel a tela."]),
    ("Lancamentos",    ORANGE,
     ["Tabela completa de transacoes.", "Filtros avancados por 14 dimensoes.",
      "Solicitacao de ajustes com fluxo de", "aprovacao. Edicao em massa."]),
    ("Aprovacoes",     PURPLE,
     ["Fila de solicitacoes pendentes.", "Aprovacao individual ou em massa.",
      "Rastreabilidade completa:", "quem, quando e o que foi alterado."]),
    ("Dashboard",      GREEN,
     ["KPIs executivos: Receita, CV, CF,", "EBITDA calculados em tempo real.",
      "Mesmos dados do DRE Gerencial.", "Visao rapida do desempenho anual."]),
    ("Painel Admin",   RGBColor(0x92, 0x40, 0x0E),
     ["Importacao de dados via planilha.", "Gerenciamento de usuarios e perfis.",
      "Exportacao do banco com filtros.", "Central de configuracao."]),
    ("Analise de IA",  RGBColor(0x1E, 0x40, 0xAF),
     ["Narrativas automaticas sobre o DRE.", "Powered by Claude (Anthropic).",
      "Contexto filtrado por sessao ativa.", "Historico de analises salvo."]),
]

positions = [
    (0.3, 1.2), (4.6, 1.2), (8.9, 1.2),
    (0.3, 4.1), (4.6, 4.1), (8.9, 4.1),
]
for i, ((lx, ty), (title, color, desc)) in enumerate(zip(positions, modules)):
    rect(s3, lx, ty, 4.0, 2.85, WHITE)
    rect(s3, lx, ty, 4.0, 0.45, color)
    txt(s3, title, lx+0.15, ty+0.08, 3.7, 0.35, 12, bold=True, color=WHITE)
    txtml(s3, desc, lx+0.15, ty+0.55, 3.7, 2.1, 10.5, color=DARK_TEXT)

# ── SLIDE 4: DRE Gerencial ───────────────────────────────────────
s4 = prs.slides.add_slide(blank)
dark_header(s4, "DRE Gerencial  -  Funcionalidades em Detalhe")

txt(s4, "FILTROS DISPONIVEIS", 0.3, 1.3, 6, 0.38, 12, bold=True, color=CYAN)
filtros = [
    "  Ano de referencia",
    "  Periodo: Mes De / Ate",
    "  Marca (cascata)",
    "  Filial (cascata da Marca)",
    "  Tag01  -  categoria gerencial",
    "  Tag02  -  subcategoria (multi)",
    "  Tag03  -  detalhe adicional (multi)",
    "  Recorrencia: Sim / Nao / Todos",
]
for j, f in enumerate(filtros):
    txt(s4, f, 0.3, 1.72+j*0.38, 5.8, 0.36, 11, color=WHITE)

txt(s4, "MODOS DE VISAO", 6.3, 1.3, 3.5, 0.38, 12, bold=True, color=CYAN)
modos = [
    ("Consolidado", "Real + Orcado + A-1 lado a lado com deltas"),
    ("Por Cenario",  "Visualizacao focada em um unico cenario"),
    ("Por Mes",      "Evolucao mensal de cada linha Tag01"),
]
for j, (m, d) in enumerate(modos):
    rect(s4, 6.3, 1.72+j*0.92, 3.3, 0.78, MID_BLUE)
    txt(s4, m, 6.5, 1.77+j*0.92, 3.0, 0.3, 11, bold=True, color=WHITE)
    txt(s4, d, 6.5, 2.08+j*0.92, 3.0, 0.35, 9.5, color=CYAN)

txt(s4, "DRILL-DOWN E EXPORT", 10.0, 1.3, 3.2, 0.38, 12, bold=True, color=CYAN)
dr = [
    "  Clique na linha para expandir",
    "  Niveis: Marca  Filial",
    "  Fornecedor  Tag02  Tag03",
    "  Ctrl+Clique soma celulas",
    "  Ordenacao A-Z ou por valor",
    "  Tela cheia (fullscreen)",
    "  Export Excel com estilos e",
    "  drill-down abertos",
]
for j, d in enumerate(dr):
    txt(s4, d, 10.0, 1.72+j*0.38, 3.2, 0.36, 10.5, color=WHITE)

rect(s4, 0.3, 5.05, 12.73, 0.06, ORANGE)
txt(s4, "LINHAS CALCULADAS AUTOMATICAS", 0.3, 5.18, 12, 0.35, 11, bold=True, color=CYAN)
calcrows = [
    ("Margem de Contribuicao",         "Receita + Custos Variaveis"),
    ("EBITDA (S/ Rateio Raiz CSC)",    "Margem + Custos Fixos"),
    ("Rateio Raiz ADM",                "Distribuicao automatica a cada 15 min"),
]
for j, (name, desc) in enumerate(calcrows):
    lx = 0.3 + j*4.24
    rect(s4, lx, 5.6, 4.0, 1.6, MID_BLUE)
    rect(s4, lx, 5.6, 4.0, 0.06, ORANGE)
    txt(s4, name, lx+0.15, 5.72, 3.75, 0.42, 11, bold=True, color=WHITE)
    txt(s4, desc, lx+0.15, 6.2, 3.75, 0.8, 10.5, color=CYAN)

# ── SLIDE 5: Fluxo de Aprovacoes ─────────────────────────────────
s5 = prs.slides.add_slide(blank)
header(s5, "Fluxo de Ajustes e Aprovacoes")

steps = [
    ("1", "USUARIO\nSOLICITA",   MID_BLUE,
     ["Abre o lancamento", "e clica em Solicitar", "Ajuste. Preenche os",
      "campos alterados e", "uma justificativa."]),
    ("2", "FILA DE\nAPROVACAO",  ORANGE,
     ["A solicitacao entra", "com status Pendente", "na guia Aprovacoes,",
      "com registro de", "data e solicitante."]),
    ("3", "APROVADOR\nANALISA",  MID_BLUE,
     ["Admin ou Aprovador", "visualiza o antes e", "depois de cada campo",
      "e decide aprovar", "ou reprovar."]),
    ("4", "APLICADO\nNO BANCO",  GREEN,
     ["Dado atualizado na", "transacao original,", "status Ajustado e",
      "justificativa salva", "no lancamento."]),
]

for j, (num, title, color, desc) in enumerate(steps):
    lx = 0.25 + j*3.2
    rect(s5, lx, 1.4, 2.9, 4.5, WHITE)
    rect(s5, lx, 1.4, 2.9, 1.05, color)
    txt(s5, num, lx+0.15, 1.45, 0.6, 0.7, 28, bold=True, color=WHITE)
    txtml(s5, title.split("\n"), lx+0.75, 1.56, 2.0, 0.8, 13, bold=True, color=WHITE)
    txtml(s5, desc, lx+0.2, 2.6, 2.55, 2.8, 11, color=DARK_TEXT)
    if j < 3:
        txt(s5, ">", lx+2.98, 2.7, 0.3, 0.5, 20, bold=True, color=MID_BLUE)

rect(s5, 0.25, 6.05, 12.83, 1.25, DARK_BLUE)
txt(s5, "TIPOS DE AJUSTE SUPORTADOS:", 0.5, 6.12, 4, 0.38, 11, bold=True, color=CYAN)
tipos = ["Data de Competencia", "Filial", "Conta Contabil", "Recorrencia", "Rateio", "Exclusao", "Multi-campo"]
for j, tp in enumerate(tipos):
    rect(s5, 0.45+j*1.78, 6.57, 1.65, 0.52, MID_BLUE)
    txt(s5, tp, 0.5+j*1.78, 6.63, 1.55, 0.4, 8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDE 6: Perfis de Acesso ─────────────────────────────────────
s6 = prs.slides.add_slide(blank)
header(s6, "Perfis de Acesso e Permissoes")

profiles = [
    ("ADMIN", ORANGE, [
        ("Acesso total a plataforma",                True),
        ("Gerencia usuarios e permissoes",           True),
        ("Solicita ajustes: Real, Orcado e A-1",    True),
        ("Aprova e rejeita solicitacoes",            True),
        ("Acessa painel administrativo",             True),
        ("Exporta banco de dados completo",          True),
        ("Restricoes granulares por marca/filial",   True),
    ]),
    ("APROVADOR", MID_BLUE, [
        ("Visualiza todos os lancamentos",           True),
        ("Solicita ajustes apenas no Real",          True),
        ("Aprova e rejeita solicitacoes",            True),
        ("Acesso ao DRE Gerencial completo",         True),
        ("Exportacao Excel",                         True),
        ("Sem acesso ao painel admin",               False),
        ("Sem ajustes em Orcado e A-1",              False),
    ]),
    ("USUARIO", GREEN, [
        ("Visualiza lancamentos (Real)",             True),
        ("Solicita ajustes apenas no Real",          True),
        ("Acesso ao DRE Gerencial",                  True),
        ("Exportacao Excel",                         True),
        ("Nao aprova solicitacoes",                  False),
        ("Sem acesso ao painel admin",               False),
        ("Sem ajustes em Orcado e A-1",              False),
    ]),
]

for j, (role, color, perms) in enumerate(profiles):
    lx = 0.3 + j*4.35
    rect(s6, lx, 1.35, 4.1, 5.85, WHITE)
    rect(s6, lx, 1.35, 4.1, 0.65, color)
    txt(s6, role, lx+0.2, 1.42, 3.7, 0.5, 18, bold=True, color=WHITE)
    for k, (perm, ok) in enumerate(perms):
        icon = "OK" if ok else "X "
        c = GREEN if ok else RGBColor(0xC0, 0x20, 0x20)
        sym = "[v]  " if ok else "[x]  "
        txt(s6, sym + perm, lx+0.2, 2.15+k*0.57, 3.75, 0.5, 11, color=c)

rect(s6, 0.3, 7.08, 12.73, 0.32, DARK_BLUE)
txt(s6, "Permissoes granulares: Admin pode restringir acesso por Marca, Filial, Conta Contabil e Tag01 por usuario.", 0.5, 7.1, 12.5, 0.28, 9.5, color=CYAN)

# ── SLIDE 7: Arquitetura Tecnica ──────────────────────────────────
s7 = prs.slides.add_slide(blank)
dark_header(s7, "Arquitetura Tecnica")

layers = [
    ("FRONTEND",     MID_BLUE,
     ["React 18 + TypeScript + Vite", "Tailwind CSS (design system)",
      "ExcelJS (export formatado)", "Recharts / ECharts (graficos)",
      "Deploy: Vercel (CDN global)"]),
    ("AUTENTICACAO", ORANGE,
     ["Firebase Authentication", "Google OAuth 2.0",
      "JWT Token", "Row Level Security (RLS)",
      "Perfis: admin, aprovador, usuario"]),
    ("BANCO DE DADOS", PURPLE,
     ["Supabase (PostgreSQL)", "RPC Functions (get_soma_tags)",
      "Materialized Views (dre_agg)", "RLS por perfil de usuario",
      "pg_cron: rateio a cada 15 min"]),
    ("INTELIGENCIA",  GREEN,
     ["Anthropic Claude Haiku", "Narrativas do DRE em linguagem natural",
      "Contexto filtrado por sessao", "Historico salvo por hash de filtro",
      "API segura via backend"]),
]

for j, (layer, color, desc) in enumerate(layers):
    lx = 0.3 + j*3.26
    rect(s7, lx, 1.35, 3.0, 3.5, HEADER_BLUE)
    rect(s7, lx, 1.35, 3.0, 0.5, color)
    txt(s7, layer, lx+0.15, 1.42, 2.75, 0.38, 11, bold=True, color=WHITE)
    txtml(s7, desc, lx+0.15, 1.98, 2.75, 2.7, 10.5, color=CYAN)

rect(s7, 0.3, 5.1, 12.73, 0.06, ORANGE)
txt(s7, "FLUXO DE DADOS", 0.3, 5.25, 4, 0.35, 11, bold=True, color=CYAN)

flow = [
    ("Planilha\nou ERP",          RGBColor(0x37, 0x41, 0x51)),
    ("Importacao\nvia Admin",     MID_BLUE),
    ("Supabase\nPostgreSQL",      PURPLE),
    ("RPC / Views\nMaterializadas", GREEN),
    ("React UI\nFiltros + DRE",   MID_BLUE),
    ("Usuario\nFinal",            ORANGE),
]
for j, (label, color) in enumerate(flow):
    lx = 0.3 + j*2.15
    rect(s7, lx, 5.7, 1.9, 1.55, color)
    txtml(s7, label.split("\n"), lx+0.1, 5.88, 1.7, 1.1, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if j < 5:
        txt(s7, ">", lx+1.95, 6.2, 0.28, 0.5, 15, bold=True, color=CYAN)

# ── SLIDE 8: Encerramento ─────────────────────────────────────────
s8 = prs.slides.add_slide(blank)
rect(s8, 0, 0, 13.33, 7.5, DARK_BLUE)
rect(s8, 0, 5.5, 13.33, 2.0, MID_BLUE)
rect(s8, 0, 5.45, 13.33, 0.08, CYAN)
rect(s8, 0.5, 1.5, 0.08, 3.5, ORANGE)
txt(s8, "RAIZ EDUCACAO", 0.75, 1.42, 10, 0.6, 13, bold=True, color=CYAN)
txt(s8, "Uma plataforma.", 0.75, 2.05, 11, 1.0, 40, bold=True, color=WHITE)
txt(s8, "Toda a inteligencia financeira.", 0.75, 2.95, 11, 1.0, 40, bold=True, color=CYAN)
txt(s8, "Do lancamento a decisao, em tempo real.", 0.75, 4.05, 10, 0.6, 16, color=RGBColor(0xB0, 0xC8, 0xE0))
txt(s8, "Acesse agora:", 0.75, 5.72, 4, 0.38, 11, color=WHITE)
txt(s8, "https://dre-raiz.vercel.app", 0.75, 6.1, 7, 0.5, 15, bold=True, color=CYAN)
txt(s8, "Raiz Educacao S.A.  |  2026  |  Confidencial", 7.0, 6.8, 6, 0.4, 9, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.RIGHT)

import os
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "DRE_Raiz_Apresentacao.pptx")
prs.save(output)
print("OK:" + output)
